import os
from pptx import Presentation
from pdf2image import convert_from_path

# Function to update the PowerPoint title and save as a PNG directly
def create_thumbnail_from_template(template_path, output_dir, video_name):
    # Load the PowerPoint template
    prs = Presentation(template_path)
    
    # Assuming the first slide contains the title we need to change
    slide = prs.slides[0]

    # Find the first title placeholder and change its text
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Assuming the title we want to change is the first text box found
            shape.text = video_name
            break

    # Ensure the output directory exists
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Save the modified slide as a temporary PowerPoint file (since we need a pptx to convert to pdf)
    temp_pptx_path = os.path.join(output_dir, "temp.pptx")
    prs.save(temp_pptx_path)

    # Save the PowerPoint file as a PDF (this is required for converting to PNG)
    temp_pdf_path = os.path.join(output_dir, "temp.pdf")
    os.system(f'libreoffice --headless --convert-to pdf "{temp_pptx_path}" --outdir "{output_dir}"')

    # Convert the first page of the PDF to PNG (as thumbnail)
    pages = convert_from_path(temp_pdf_path, dpi=300)
    thumbnail_path = os.path.join(output_dir, f"{video_name}.png")

    # Save the first page as PNG
    pages[0].save(thumbnail_path, 'PNG')

    # Cleanup temporary files
    os.remove(temp_pptx_path)  # Remove the temporary pptx file
    os.remove(temp_pdf_path)   # Remove the temporary pdf file

    return thumbnail_path

# Function to process input file and generate thumbnails
def process_video_thumbnails(input_file, template_path, output_dir):
    with open(input_file, 'r') as f:
        lines = f.readlines()

    # Process lines in pairs (name, url)
    for i in range(0, len(lines), 2):
        video_name = lines[i].strip()  # Video name
        embed_url = lines[i + 1].strip()  # Embed URL (not used here for thumbnails)

        print(f"Generating thumbnail for: {video_name}")
        
        # Create thumbnail for this video
        thumbnail_path = create_thumbnail_from_template(template_path, output_dir, video_name)
        print(f"Thumbnail saved to: {thumbnail_path}")

# Example usage
if __name__ == "__main__":
    input_file = 'video_links.txt'  # Path to your input file with video names and URLs
    template_path = 'RMIT_lecture_video_thumb_template.pptx'  # Path to the PowerPoint template
    output_dir = './thumbnails'  # Directory to save the generated thumbnails

    process_video_thumbnails(input_file, template_path, output_dir)
