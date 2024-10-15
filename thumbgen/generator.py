import os
import tempfile
import subprocess
import argparse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pdf2image import convert_from_path

def create_thumbnail_from_template(template_path, output_dir, video_name, change_text_box_name):
    try:
        safe_video_name = "".join(c for c in video_name if c.isalnum() or c in (' ', '_', '-')).rstrip()

        prs = Presentation(template_path)
        slide = prs.slides[0]
        white_color = RGBColor(255, 255, 255)
        font_size = Pt(16)

        dynamic_text_box = None
        for shape in slide.shapes:
            if shape.has_text_frame and change_text_box_name in shape.text:
                dynamic_text_box = shape
                break

        if not dynamic_text_box:
            print(f"Error: Could not find the '{change_text_box_name}' text box on the slide for '{video_name}'.")
            return None

        dynamic_text_box.text = video_name
        for paragraph in dynamic_text_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = white_color
                run.font.size = font_size

        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx', dir=output_dir) as tmp_pptx_file:
            temp_pptx_path = tmp_pptx_file.name
        prs.save(temp_pptx_path)
        prs = None

        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to',
            'pdf',
            temp_pptx_path,
            '--outdir',
            output_dir
        ]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

        if result.returncode != 0 or not os.path.exists(os.path.splitext(temp_pptx_path)[0] + ".pdf"):
            print(f"LibreOffice conversion to PDF failed for '{video_name}'.")
            print(f"Command: {' '.join(cmd)}")
            print(f"Return Code: {result.returncode}")
            print(f"Stdout: {result.stdout}")
            print(f"Stderr: {result.stderr}")
            os.remove(temp_pptx_path)
            return None

        temp_pdf_path = os.path.splitext(temp_pptx_path)[0] + ".pdf"

        pages = convert_from_path(temp_pdf_path, dpi=300)
        thumbnail_path = os.path.join(output_dir, f"{safe_video_name}.png")
        pages[0].save(thumbnail_path, 'PNG')

        os.remove(temp_pptx_path)
        os.remove(temp_pdf_path)

        return thumbnail_path
    except Exception as e:
        print(f"Exception occurred while processing '{video_name}': {e}")
        return None

def process_video_thumbnails(input_file, template_path, output_dir, change_text_box_name):
    with open(input_file, 'r') as f:
        lines = f.readlines()

    for line in lines:
        video_name = line.strip()

        print(f"Generating thumbnail for: {video_name}")

        thumbnail_path = create_thumbnail_from_template(template_path, output_dir, video_name, change_text_box_name)
        if thumbnail_path:
            print(f"Thumbnail for '{video_name}' saved to: {thumbnail_path}")
        else:
            print(f"Failed to create thumbnail for '{video_name}'.")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Generate thumbnails for videos using a PowerPoint template.')

    parser.add_argument('--input_file', type=str, required=True, help='Path to the input file containing video names.')
    parser.add_argument('--template_path', type=str, required=True, help='Path to the PowerPoint (.pptx) template file.')
    parser.add_argument('--output_dir', type=str, required=True, help='Directory where thumbnails will be saved.')
    parser.add_argument('--text_box_name', type=str, required=True, help='The name or text in the text box to be replaced in the template.')

    args = parser.parse_args()

    process_video_thumbnails(args.input_file, args.template_path, args.output_dir, args.text_box_name)
