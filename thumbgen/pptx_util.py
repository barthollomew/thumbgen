import os
import tempfile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from thumbgen.pdf_util import convert_pptx_to_png
from thumbgen.compress_util import compress_image


def create_thumbnail_from_template(
    template_path,
    output_dir,
    video_name,
    change_text_box_name,
    text_color="white",
    font_size=16,
):
    """
    Generate a PNG thumbnail from a PPTX template, then compress it if >3MB.

    Returns the path to the PNG or None on failure.
    """
    try:
        # sanitize file-name
        safe_video_name = "".join(
            c for c in video_name if c.isalnum() or c in (" ", "_", "-")
        ).rstrip()

        # load template and modify text
        prs = Presentation(template_path)
        slide = prs.slides[0]

        rgb = (
            RGBColor(255, 255, 255)
            if text_color.lower() == "white"
            else RGBColor(0, 0, 0)
        )
        font_size_pt = Pt(font_size)

        dynamic_text_box = None
        for shape in slide.shapes:
            if shape.has_text_frame and change_text_box_name in shape.text:
                dynamic_text_box = shape
                break

        if not dynamic_text_box:
            print(
                f"Error: Could not find '{change_text_box_name}' in template for '{video_name}'"
            )
            return None

        dynamic_text_box.text = video_name
        for paragraph in dynamic_text_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = rgb
                run.font.size = font_size_pt

        # ensure output directory
        os.makedirs(output_dir, exist_ok=True)

        # save to a temp PPTX
        with tempfile.NamedTemporaryFile(
            delete=False, suffix=".pptx", dir=output_dir
        ) as tmp_ppt:
            temp_pptx_path = tmp_ppt.name
        prs.save(temp_pptx_path)
        prs = None

        # convert to PNG
        thumbnail_path = convert_pptx_to_png(
            temp_pptx_path, output_dir, safe_video_name
        )

        # compress if needed
        if thumbnail_path:
            try:
                compress_image(thumbnail_path)
            except Exception as e:
                print(f"Compression error on {thumbnail_path}: {e}")

        return thumbnail_path

    except Exception as e:
        print(f"Exception occurred while processing '{video_name}': {e}")
        return None
